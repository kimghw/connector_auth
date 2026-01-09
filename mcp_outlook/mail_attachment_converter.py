"""
TXT Converter - 첨부파일 텍스트 변환

지원 포맷:
    - PDF → TXT (pdfplumber)
    - DOCX/DOC → TXT (python-docx)
    - HWP/HWPX → HTML → TXT (pyhwpx)
    - XLSX/XLS → TXT (openpyxl)
    - PPTX → TXT (python-pptx)

Classes:
    - FileConverter (ABC): 변환기 추상 인터페이스
    - PdfConverter: PDF 변환
    - WordConverter: Word 문서 변환
    - HwpConverter: 한글 문서 변환
    - ExcelConverter: Excel 변환
    - PowerPointConverter: PowerPoint 변환
    - ConversionPipeline: 변환기 라우팅
"""

import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional, List, Tuple
from io import BytesIO


class FileConverter(ABC):
    """
    파일 변환 추상 인터페이스

    각 파일 포맷별 변환기가 구현해야 하는 메서드 정의
    """

    @abstractmethod
    def supports(self, extension: str) -> bool:
        """
        해당 확장자 지원 여부

        Args:
            extension: 파일 확장자 (예: ".pdf")

        Returns:
            지원 여부
        """
        pass

    @abstractmethod
    def convert(self, content: bytes, filename: str) -> str:
        """
        파일을 텍스트로 변환

        Args:
            content: 파일 내용 (bytes)
            filename: 파일명

        Returns:
            변환된 텍스트

        Raises:
            Exception: 변환 실패 시
        """
        pass

    @property
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        """지원하는 확장자 목록"""
        pass


class PdfConverter(FileConverter):
    """
    PDF → TXT 변환

    pdfplumber 사용
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".pdf"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def convert(self, content: bytes, filename: str) -> str:
        """
        PDF 파일을 텍스트로 변환

        Args:
            content: PDF 파일 내용
            filename: 파일명

        Returns:
            추출된 텍스트
        """
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber가 필요합니다: pip install pdfplumber")

        text_parts = []

        with pdfplumber.open(BytesIO(content)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {i} ---")
                    text_parts.append(page_text)

        return "\n\n".join(text_parts)


class WordConverter(FileConverter):
    """
    Word 문서 (DOCX) → TXT 변환

    python-docx 사용
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".docx", ".doc"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def convert(self, content: bytes, filename: str) -> str:
        """
        Word 문서를 텍스트로 변환

        Args:
            content: Word 파일 내용
            filename: 파일명

        Returns:
            추출된 텍스트
        """
        ext = Path(filename).suffix.lower()

        if ext == ".doc":
            # .doc 파일은 antiword 또는 다른 방법 필요
            # 여기서는 기본적으로 지원하지 않음
            raise NotImplementedError(".doc 파일은 현재 지원되지 않습니다. .docx로 변환 후 시도하세요.")

        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx가 필요합니다: pip install python-docx")

        doc = Document(BytesIO(content))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # 테이블 내용도 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))

        return "\n".join(text_parts)


class HwpConverter(FileConverter):
    """
    한글 문서 (HWP/HWPX) → TXT 변환

    pyhwpx 사용하여 HTML로 변환 후 텍스트 추출
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".hwp", ".hwpx"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def _html_to_text(self, html_content: str) -> str:
        """
        HTML을 텍스트로 변환

        Args:
            html_content: HTML 문자열

        Returns:
            텍스트
        """
        # HTML 태그 제거
        text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<br\s*/?>', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'<p[^>]*>', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'</p>', '', text, flags=re.IGNORECASE)
        text = re.sub(r'<[^>]+>', '', text)

        # HTML 엔티티 변환
        text = re.sub(r'&nbsp;', ' ', text)
        text = re.sub(r'&lt;', '<', text)
        text = re.sub(r'&gt;', '>', text)
        text = re.sub(r'&amp;', '&', text)
        text = re.sub(r'&quot;', '"', text)
        text = re.sub(r'&#39;', "'", text)

        # 연속된 공백/줄바꿈 정리
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)

        return text.strip()

    def convert(self, content: bytes, filename: str) -> str:
        """
        HWP/HWPX 파일을 텍스트로 변환

        Args:
            content: HWP 파일 내용
            filename: 파일명

        Returns:
            추출된 텍스트
        """
        try:
            import pyhwpx
        except ImportError:
            raise ImportError("pyhwpx가 필요합니다: pip install pyhwpx")

        try:
            # pyhwpx로 HWP 파일 열기
            hwp = pyhwpx.Hwp()
            hwp.open(BytesIO(content))

            # HTML로 변환
            html_content = hwp.to_html()

            # HTML → 텍스트
            text = self._html_to_text(html_content)

            hwp.close()
            return text

        except Exception as e:
            # pyhwpx 실패 시 대안: olefile 사용
            try:
                return self._convert_with_olefile(content)
            except Exception:
                raise Exception(f"HWP 변환 실패: {e}")

    def _convert_with_olefile(self, content: bytes) -> str:
        """
        olefile을 사용한 대안 변환

        Args:
            content: HWP 파일 내용

        Returns:
            추출된 텍스트
        """
        try:
            import olefile
            import zlib
        except ImportError:
            raise ImportError("olefile이 필요합니다: pip install olefile")

        ole = olefile.OleFileIO(BytesIO(content))

        # HWP 본문 스트림 찾기
        if ole.exists("BodyText/Section0"):
            data = ole.openstream("BodyText/Section0").read()

            # 압축 해제 시도
            try:
                data = zlib.decompress(data, -15)
            except zlib.error:
                pass

            # 텍스트 추출 (간단한 방법)
            text = data.decode("utf-16", errors="ignore")
            # 제어 문자 제거
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

            ole.close()
            return text

        ole.close()
        raise Exception("HWP 본문을 찾을 수 없습니다")


class ExcelConverter(FileConverter):
    """
    Excel (XLSX) → TXT 변환

    openpyxl 사용
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".xlsx", ".xls"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def convert(self, content: bytes, filename: str) -> str:
        """
        Excel 파일을 텍스트로 변환

        Args:
            content: Excel 파일 내용
            filename: 파일명

        Returns:
            추출된 텍스트
        """
        ext = Path(filename).suffix.lower()

        if ext == ".xls":
            raise NotImplementedError(".xls 파일은 현재 지원되지 않습니다. .xlsx로 변환 후 시도하세요.")

        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl이 필요합니다: pip install openpyxl")

        wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows(values_only=True):
                row_values = []
                for cell in row:
                    if cell is not None:
                        row_values.append(str(cell))
                if row_values:
                    text_parts.append("\t".join(row_values))

            text_parts.append("")

        wb.close()
        return "\n".join(text_parts)


class PowerPointConverter(FileConverter):
    """
    PowerPoint (PPTX) → TXT 변환

    python-pptx 사용
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def convert(self, content: bytes, filename: str) -> str:
        """
        PowerPoint 파일을 텍스트로 변환

        Args:
            content: PowerPoint 파일 내용
            filename: 파일명

        Returns:
            추출된 텍스트
        """
        ext = Path(filename).suffix.lower()

        if ext == ".ppt":
            raise NotImplementedError(".ppt 파일은 현재 지원되지 않습니다. .pptx로 변환 후 시도하세요.")

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx가 필요합니다: pip install python-pptx")

        prs = Presentation(BytesIO(content))
        text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {i} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

            text_parts.append("")

        return "\n".join(text_parts)


class PlainTextConverter(FileConverter):
    """
    텍스트 파일 처리 (변환 없이 디코딩만)
    """

    @property
    def supported_extensions(self) -> List[str]:
        return [".txt", ".csv", ".json", ".xml", ".html", ".htm", ".md", ".log"]

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def convert(self, content: bytes, filename: str) -> str:
        """
        텍스트 파일 디코딩

        Args:
            content: 파일 내용
            filename: 파일명

        Returns:
            디코딩된 텍스트
        """
        # 인코딩 감지 시도
        encodings = ["utf-8", "cp949", "euc-kr", "latin-1"]

        for encoding in encodings:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        # 모든 인코딩 실패 시 errors='replace' 사용
        return content.decode("utf-8", errors="replace")


class ConversionPipeline:
    """
    파일 변환 파이프라인

    확장자에 따라 적절한 변환기 선택 및 실행
    """

    def __init__(self):
        """초기화 - 모든 변환기 등록"""
        self.converters: List[FileConverter] = [
            PdfConverter(),
            WordConverter(),
            HwpConverter(),
            ExcelConverter(),
            PowerPointConverter(),
            PlainTextConverter(),
        ]

    def get_supported_extensions(self) -> List[str]:
        """
        지원하는 모든 확장자 목록

        Returns:
            확장자 목록
        """
        extensions = []
        for converter in self.converters:
            extensions.extend(converter.supported_extensions)
        return extensions

    def can_convert(self, filename: str) -> bool:
        """
        파일 변환 가능 여부 확인

        Args:
            filename: 파일명

        Returns:
            변환 가능 여부
        """
        ext = Path(filename).suffix.lower()
        return any(c.supports(ext) for c in self.converters)

    def get_converter(self, filename: str) -> Optional[FileConverter]:
        """
        파일에 맞는 변환기 반환

        Args:
            filename: 파일명

        Returns:
            변환기 또는 None
        """
        ext = Path(filename).suffix.lower()
        for converter in self.converters:
            if converter.supports(ext):
                return converter
        return None

    def convert(self, content: bytes, filename: str) -> Tuple[Optional[str], Optional[str]]:
        """
        파일을 텍스트로 변환

        Args:
            content: 파일 내용 (bytes)
            filename: 파일명

        Returns:
            (변환된 텍스트, 에러 메시지) 튜플
            성공 시: (text, None)
            실패 시: (None, error_message)
        """
        converter = self.get_converter(filename)

        if not converter:
            return None, f"지원하지 않는 파일 형식: {Path(filename).suffix}"

        try:
            text = converter.convert(content, filename)
            return text, None
        except ImportError as e:
            return None, f"필요한 라이브러리 없음: {e}"
        except NotImplementedError as e:
            return None, str(e)
        except Exception as e:
            return None, f"변환 실패: {e}"

    def convert_to_txt_filename(self, original_filename: str) -> str:
        """
        원본 파일명을 .txt 확장자로 변경

        Args:
            original_filename: 원본 파일명

        Returns:
            .txt 확장자 파일명
        """
        stem = Path(original_filename).stem
        return f"{stem}.txt"


# 편의를 위한 싱글톤 인스턴스
_pipeline_instance: Optional[ConversionPipeline] = None


def get_conversion_pipeline() -> ConversionPipeline:
    """
    ConversionPipeline 싱글톤 인스턴스 반환

    Returns:
        ConversionPipeline 인스턴스
    """
    global _pipeline_instance
    if _pipeline_instance is None:
        _pipeline_instance = ConversionPipeline()
    return _pipeline_instance
